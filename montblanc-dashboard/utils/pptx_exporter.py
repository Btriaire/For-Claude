import os
import tempfile
from datetime import datetime
from typing import Dict, List, Optional

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


GEHC_PALETTE = [
    "#6400A0", "#7845B3", "#66BDCC", "#31C18A",
    "#F8DA65", "#F2886E", "#B02DFF", "#CB73FF",
]

GEHC_FONT = "Source Sans Pro"


def _rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _format_value(value: float, fmt: str) -> str:
    if fmt == "currency":
        if value >= 1_000_000:
            return f"${value / 1_000_000:.1f}M"
        if value >= 1_000:
            return f"${value / 1_000:.1f}K"
        return f"${value:,.0f}"
    if fmt == "percentage":
        return f"{value:.1f}%"
    if value >= 1_000_000:
        return f"{value / 1_000_000:.1f}M"
    if value >= 1_000:
        return f"{value / 1_000:.1f}K"
    return f"{value:,.0f}"


PALETTE = [
    "#5B2D8E", "#C9A227", "#6DBF8B", "#E8735A",
    "#9B59B6", "#D4A017", "#52BE80", "#E07B5A",
]

CARD_COLORS = ["#5B2D8E", "#C9A227", "#6DBF8B", "#E8735A",
               "#7D3C98", "#B7860B", "#58D68D", "#D35400"]


class PPTXExporter:
    W = Inches(13.33)
    H = Inches(7.5)

    def __init__(self, config: Dict):
        self.config = config
        colors = config.get("colors", {})
        # Support both new c1-c8 keys and legacy primary/secondary keys
        self.palette = [
            colors.get(f"c{i}", GEHC_PALETTE[i - 1]) for i in range(1, 9)
        ]
        self.purple = self.palette[0]
        self.yellow = self.palette[4]
        self.green  = self.palette[3]
        self.orange = self.palette[5]

    def generate(self, data: List[Dict], period: str,
                 comments: str, key_events: str) -> str:
        prs = Presentation()
        prs.slide_width  = self.W
        prs.slide_height = self.H

        self._title_slide(prs, period)
        if data:
            self._kpi_overview_slide(prs, data, period)
        chart_kpis = [d for d in data if d["breakdown"]]
        for i in range(0, len(chart_kpis), 2):
            self._charts_slide(prs, chart_kpis[i:i + 2], period)
        if key_events or comments:
            self._events_slide(prs, key_events, comments, period)

        out = tempfile.mktemp(suffix=".pptx")
        prs.save(out)
        return out

    # ── helpers ──────────────────────────────────────────────────────────────

    def _blank(self, prs):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = _rgb("#FFFFFF")
        return slide

    def _rect(self, slide, x, y, w, h, fill_hex):
        s = slide.shapes.add_shape(1, x, y, w, h)
        s.fill.solid()
        s.fill.fore_color.rgb = _rgb(fill_hex)
        s.line.fill.background()
        return s

    def _tb(self, slide, x, y, w, h, text, size,
            bold=False, color="#1A1A1A", align=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.name  = GEHC_FONT
        r.font.size  = Pt(size)
        r.font.bold  = bold
        r.font.color.rgb = _rgb(color)

    def _header(self, slide, prs, title, period):
        self._rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.9), self.purple)
        self._rect(slide, Inches(0), Inches(0.9), prs.slide_width, Inches(0.06), self.yellow)
        self._tb(slide, Inches(0.35), Inches(0.1), Inches(9), Inches(0.75),
                 title, 22, bold=True, color="#FFFFFF")
        self._tb(slide, Inches(9.8), Inches(0.2), Inches(3.2), Inches(0.5),
                 period, 14, color="#FFFFFF", align=PP_ALIGN.RIGHT)

    # ── slides ────────────────────────────────────────────────────────────────

    def _title_slide(self, prs, period):
        slide = self._blank(prs)
        self._rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.15), self.purple)
        self._rect(slide, Inches(0), prs.slide_height - Inches(0.15),
                   prs.slide_width, Inches(0.15), self.purple)
        self._rect(slide, Inches(0.5), Inches(2.3), Inches(0.12), Inches(3.0), self.purple)
        self._rect(slide, Inches(0.5), Inches(5.3), Inches(6), Inches(0.06), self.yellow)
        self._tb(slide, Inches(0.8), Inches(2.5), Inches(10), Inches(1.3),
                 "GE HealthCare — INT STO Enterprise Solutions", 18, bold=True, color="#7F7F7F")
        self._tb(slide, Inches(0.8), Inches(3.1), Inches(10), Inches(1.3),
                 "Performance Dashboard", 28, bold=True, color="#000000")
        self._tb(slide, Inches(0.8), Inches(4.1), Inches(8), Inches(0.65),
                 period, 24, color=self.purple)
        self._tb(slide, Inches(0.8), Inches(4.9), Inches(8), Inches(0.45),
                 f"Generated on {datetime.now().strftime('%B %d, %Y')}", 10, color="#7F7F7F")

    def _kpi_overview_slide(self, prs, data, period):
        slide = self._blank(prs)
        self._header(slide, prs, "KPI Overview", period)
        card_w, card_h, gap = Inches(2.8), Inches(1.55), Inches(0.28)
        for i, kpi in enumerate(data[:8]):
            row, col = divmod(i, 4)
            x = Inches(0.45) + col * (card_w + gap)
            y = Inches(1.25) + row * (card_h + gap)
            color = self.palette[i % len(self.palette)]
            self._rect(slide, x, y, card_w, card_h, color)
            self._tb(slide, x + Inches(0.15), y + Inches(0.12),
                     card_w - Inches(0.3), Inches(0.38),
                     kpi["label"].upper(), 10, bold=True, color="#FFFFFF")
            self._tb(slide, x + Inches(0.15), y + Inches(0.52),
                     card_w - Inches(0.3), Inches(0.85),
                     _format_value(kpi["total"], kpi["format"]),
                     28, bold=True, color="#FFFFFF")

    def _charts_slide(self, prs, kpis, period):
        slide = self._blank(prs)
        title = "  ·  ".join(k["label"] for k in kpis)
        self._header(slide, prs, title, period)
        positions = [(Inches(0.4), Inches(1.1)), (Inches(6.9), Inches(1.1))]
        for i, kpi in enumerate(kpis[:2]):
            x, y = positions[i]
            if kpi.get("aggregation") == "histogram" or kpi.get("chart_type") == "bar":
                self._add_column_chart(slide, kpi, x, y, Inches(5.9), Inches(5.8))
            else:
                self._add_donut_chart(slide, kpi, x, y, Inches(5.9), Inches(5.8))

    def _events_slide(self, prs, key_events, comments, period):
        slide = self._blank(prs)
        self._header(slide, prs, "Key Events & Commentary", period)
        y = Inches(1.15)
        if key_events:
            self._tb(slide, Inches(0.5), y, Inches(12.3), Inches(0.4),
                     "KEY EVENTS", 13, bold=True, color=self.purple)
            y += Inches(0.45)
            self._rect(slide, Inches(0.5), y, Inches(0.06), Inches(2.0), self.yellow)
            self._tb(slide, Inches(0.75), y, Inches(12.0), Inches(2.0),
                     key_events, 13, color="#1A1A1A")
            y += Inches(2.3)
        if comments:
            self._tb(slide, Inches(0.5), y, Inches(12.3), Inches(0.4),
                     "COMMENTARY", 13, bold=True, color=self.purple)
            y += Inches(0.45)
            self._rect(slide, Inches(0.5), y, Inches(0.06), Inches(1.8), self.green)
            self._tb(slide, Inches(0.75), y, Inches(12.0), Inches(1.8),
                     comments, 13, color="#1A1A1A")

    def _add_column_chart(self, slide, kpi, x, y, w, h):
        if not kpi["breakdown"]:
            return

        labels = [b["label"] for b in kpi["breakdown"]]
        values = [b["value"] for b in kpi["breakdown"]]
        total  = sum(values)

        cd = ChartData()
        cd.categories = labels
        cd.add_series(kpi["label"], values)

        chart_frame = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, w, h, cd)
        chart = chart_frame.chart

        chart.has_title = True
        chart.chart_title.text_frame.text = (
            f"{kpi['label']}  —  {_format_value(total, kpi['format'])}"
        )
        chart.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(13)
        chart.chart_title.text_frame.paragraphs[0].runs[0].font.bold = True
        chart.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = _rgb("#1A1A1A")

        is_histogram = kpi.get("aggregation") == "histogram"
        series = chart.series[0]
        for i, point in enumerate(series.points):
            color_hex = self.purple if is_histogram else self.palette[i % len(self.palette)]
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = _rgb(color_hex)

        chart.has_legend = not is_histogram
        if chart.has_legend:
            chart.legend.include_in_layout = False

    def _add_donut_chart(self, slide, kpi, x, y, w, h):
        if not kpi["breakdown"]:
            return

        labels = [b["label"] for b in kpi["breakdown"]]
        values = [b["value"] for b in kpi["breakdown"]]
        total  = sum(values)

        # Native python-pptx doughnut chart
        cd = ChartData()
        cd.categories = labels
        cd.add_series(kpi["label"], values)

        chart_frame = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, x, y, w, h, cd)
        chart = chart_frame.chart

        # Title
        chart.has_title = True
        chart.chart_title.text_frame.text = (
            f"{kpi['label']}  —  {_format_value(total, kpi['format'])}"
        )
        chart.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(13)
        chart.chart_title.text_frame.paragraphs[0].runs[0].font.bold = True
        chart.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = _rgb("#1A1A1A")

        # Color each slice
        series = chart.series[0]
        for i, point in enumerate(series.points):
            color_hex = self.palette[i % len(self.palette)]
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = _rgb(color_hex)

        # Legend
        chart.has_legend = True
        chart.legend.include_in_layout = False
